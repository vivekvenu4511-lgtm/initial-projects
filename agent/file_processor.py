# ============================================================
#  File Processing Utilities
# ============================================================

from __future__ import annotations

import csv
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))
import config as cfg


def _safe_stem(name: str, fallback: str = "output") -> str:
    """Returns a safe file stem (name without extension). Uses a fallback name if necessary."""
    stem = Path(name or fallback).stem.strip() or fallback
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem)
    return stem[:80].strip("._-") or fallback


def build_output_path(output_name: str, ext: str) -> Path:
    """Builds the path for the output file based on the specified output name and file extension."""
    ext = ext if ext.startswith(".") else f".{ext}"
    stem = _safe_stem(output_name, fallback="output")
    if not stem.lower().endswith(ext.lower()):
        stem = f"{stem}{ext}"
    out = Path(cfg.GENERATED_DIR) / stem
    out.parent.mkdir(parents=True, exist_ok=True)
    if out.exists():
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        out = out.with_name(f"{out.stem}_{ts}{out.suffix}")
    return out


def file_metadata(filepath: str | Path) -> dict:
    """Returns metadata for a given file, such as size and modification time."""
    path = Path(filepath)
    return {
        "name": path.name,
        "path": str(path),
        "size": path.stat().st_size if path.exists() else 0,
        "filetype": path.suffix.lower(),
        "modified": datetime.fromtimestamp(path.stat().st_mtime).isoformat() if path.exists() else None,
    }


def read_excel(filepath: str) -> dict:
    """Reads the contents of an Excel file and returns it as a dictionary."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        result = {"filename": Path(filepath).name, "sheets": {}}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                clean = [str(c) if c is not None else "" for c in row]
                if any(clean):
                    rows.append(clean)
            result["sheets"][sheet_name] = {
                "rows": rows[:200],
                "total_rows": ws.max_row or 0,
                "total_cols": ws.max_column or 0,
                "truncated": (ws.max_row or 0) > 200,
            }
        return result
    except Exception as e:
        return {"error": str(e)}


def read_csv_file(filepath: str) -> dict:
    """Reads a CSV file and returns its contents as a dictionary."""
    try:
        rows = []
        with open(filepath, "r", encoding="utf-8", errors="replace", newline="") as f:
            for row in csv.reader(f):
                if any(cell.strip() for cell in row):
                    rows.append(row)
        max_cols = max((len(r) for r in rows), default=0)
        return {"filename": Path(filepath).name, "sheets": {"CSV": {"rows": rows[:200], "total_rows": len(rows), "total_cols": max_cols, "truncated": len(rows) > 200}}}
    except Exception as e:
        return {"error": str(e)}


def excel_to_text(filepath: str) -> str:
    """Converts an Excel file into plain text format."""
    data = read_csv_file(filepath) if Path(filepath).suffix.lower() == ".csv" else read_excel(filepath)
    if "error" in data:
        return f"Error reading tabular file: {data['error']}"
    lines = [f"Tabular file: {data['filename']}"]
    for sname, sdata in data["sheets"].items():
        lines.append(f"\n── Sheet: {sname} ──")
        lines.append(f"Size: {sdata['total_rows']} rows × {sdata['total_cols']} columns")
        if sdata["truncated"]:
            lines.append("(showing first 200 rows)")
        for row in sdata["rows"]:
            lines.append(" | ".join(str(v) for v in row))
    return "\n".join(lines)


def compare_excel_text(path1: str, path2: str) -> str:
    """Compares two Excel files and returns the differences as text."""
    d1 = read_csv_file(path1) if Path(path1).suffix.lower() == ".csv" else read_excel(path1)
    d2 = read_csv_file(path2) if Path(path2).suffix.lower() == ".csv" else read_excel(path2)
    lines = ["=== COMPARING TABULAR FILES ===", f"File A: {d1.get('filename', Path(path1).name)}", f"File B: {d2.get('filename', Path(path2).name)}", ""]
    sheets1 = set(d1.get("sheets", {}).keys())
    sheets2 = set(d2.get("sheets", {}).keys())
    if sheets1 - sheets2:
        lines.append(f"Sheets only in A: {', '.join(sorted(sheets1 - sheets2))}")
    if sheets2 - sheets1:
        lines.append(f"Sheets only in B: {', '.join(sorted(sheets2 - sheets1))}")
    for s in sorted(sheets1 & sheets2):
        s1 = d1["sheets"][s]
        s2 = d2["sheets"][s]
        lines.append(f"\nSheet: {s}")
        lines.append(f"  A: {s1['total_rows']}r × {s1['total_cols']}c | B: {s2['total_rows']}r × {s2['total_cols']}c")
        max_rows = min(max(len(s1["rows"]), len(s2["rows"])), 50)
        for i in range(max_rows):
            r1 = s1["rows"][i] if i < len(s1["rows"]) else []
            r2 = s2["rows"][i] if i < len(s2["rows"]) else []
            if r1 != r2:
                lines.append(f"  Row {i+1} differs:")
                lines.append(f"    A: {' | '.join(map(str, r1))}")
                lines.append(f"    B: {' | '.join(map(str, r2))}")
    return "\n".join(lines)


def _normalize_rows(data: Any) -> list[dict]:
    """Normalizes the rows of a dataset to ensure consistent formatting and structure."""
    if isinstance(data, dict):
        data = [data]
    if not isinstance(data, list):
        return [{"Value": str(data)}]
    rows = []
    for idx, item in enumerate(data, 1):
        if isinstance(item, dict):
            rows.append({k: (json.dumps(v, ensure_ascii=False) if isinstance(v, (dict, list)) else v) for k, v in item.items()})
        else:
            rows.append({"Item": item, "Index": idx})
    return rows


def generate_excel(data: list[dict] | dict | list, output_name: str, sheet_name: str = "Data") -> str:
    """Generates an Excel file from a list or dictionary of data, with an optional sheet name."""
    import openpyxl
    from openpyxl.styles import Alignment, Font, PatternFill
    rows = _normalize_rows(data)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = _safe_stem(sheet_name, fallback="Data")[:31]
    headers = list(rows[0].keys()) if rows else ["Value"]
    for ci, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=ci, value=header)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="0F766E")
        cell.alignment = Alignment(horizontal="center", vertical="center")
    for ri, row in enumerate(rows, 2):
        for ci, header in enumerate(headers, 1):
            value = row.get(header, "")
            if isinstance(value, (dict, list)):
                value = json.dumps(value, ensure_ascii=False)
            ws.cell(row=ri, column=ci, value=value)
    ws.freeze_panes = "A2"
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)
    out = build_output_path(output_name, ".xlsx")
    wb.save(str(out))
    return str(out)


def read_word(filepath: str) -> str:
    """Reads the contents of a Word document and returns it as text."""
    try:
        import docx
        doc = docx.Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip()) or "(empty document)"
    except Exception as e:
        return f"Error reading Word: {e}"


def generate_word(title: str, sections: list[dict], output_name: str) -> str:
    """Generates a Word document with specified title, sections, and output name."""
    import docx
    from docx.shared import RGBColor
    doc = docx.Document()
    heading = doc.add_heading(title or "Untitled Report", 0)
    if heading.runs:
        heading.runs[0].font.color.rgb = RGBColor(15, 118, 110)
    doc.add_paragraph(f"Generated by {cfg.AGENT_NAME} on {datetime.now().strftime('%B %d, %Y %H:%M')}")
    doc.add_paragraph("")
    for sec in sections or []:
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=1)
        if sec.get("content"):
            doc.add_paragraph(str(sec["content"]))
    out = build_output_path(output_name, ".docx")
    doc.save(str(out))
    return str(out)


def generate_text(content: str, output_name: str, markdown: bool = False) -> str:
    """Generates plain text content from a provided string, with an option for markdown formatting."""
    out = build_output_path(output_name, ".md" if markdown else ".txt")
    Path(out).write_text(content, encoding="utf-8")
    return str(out)


def read_ppt(filepath: str) -> str:
    """Reads the contents of a PowerPoint presentation and returns it as text."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    texts.append(shape.text.strip())
            slides.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(slides) or "(no text content)"
    except Exception as e:
        return f"Error reading PPT: {e}"


def generate_ppt(title: str, slides: list[dict], output_name: str) -> str:
    """Generates a PowerPoint presentation from the provided title and slide data."""
    from pptx import Presentation
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title or "Presentation"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"Generated by {cfg.AGENT_NAME} · {datetime.now().strftime('%B %d, %Y')}"
    for slide in slides or []:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = str(slide.get("title") or "Slide")
        body = sl.placeholders[1].text_frame
        body.clear()
        bullets = slide.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        if not bullets:
            bullets = ["Content to be added."]
        for idx, bullet in enumerate(bullets):
            p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.level = 0
        notes = str(slide.get("notes") or "").strip()
        if notes:
            sl.notes_slide.notes_text_frame.text = notes
    out = build_output_path(output_name, ".pptx")
    prs.save(str(out))
    return str(out)


def read_pdf(filepath: str) -> str:
    """Reads the contents of a PDF file and returns it as text."""
    try:
        import pypdf
        reader = pypdf.PdfReader(filepath)
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        return f"Error reading PDF: {e}"


def read_file_for_llm(filepath: str) -> str:
    """Reads a file and prepares it for processing by a large language model (LLM)."""
    ext = Path(filepath).suffix.lower()
    if ext in (".xlsx", ".xls", ".csv"):
        return excel_to_text(filepath)
    if ext == ".docx":
        return read_word(filepath)
    if ext in (".pptx", ".ppt"):
        return read_ppt(filepath)
    if ext == ".pdf":
        return read_pdf(filepath)
    if ext in (".txt", ".md", ".json", ".py", ".js", ".html", ".css", ".sql"):
        try:
            return Path(filepath).read_text(encoding="utf-8", errors="replace")
        except Exception as e:
            return f"Error reading file: {e}"
    return f"Unsupported file type: {ext}"
